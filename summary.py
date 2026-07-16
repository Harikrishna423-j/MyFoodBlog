import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_summary_slide():
    # 1. Initialize Presentation with Widescreen (16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use a blank layout (usually layout index 6 is blank)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Define Theme Colors (Warm Food Blog Palette)
    BG_COLOR = RGBColor(250, 250, 247)       # Warm Off-White
    TEXT_DARK = RGBColor(15, 23, 42)         # Slate 900
    TEXT_MUTED = RGBColor(71, 85, 105)       # Slate 600
    ACCENT_ORANGE = RGBColor(234, 88, 12)    # Orange 600
    CARD_BG = RGBColor(255, 255, 255)        # Pure White
    BORDER_COLOR = RGBColor(226, 232, 240)   # Slate 200

    # 2. Add Background Shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background() # No border

    # 3. Add Header Section (Title + Subtitle)
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    p_title = tf.paragraphs[0]
    p_title.text = "MyFoodBlog Project Overview"
    p_title.font.name = "Georgia"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_DARK
    p_title.space_after = Pt(4)

    p_sub = tf.add_paragraph()
    p_sub.text = "Full-Stack Food Blog with Django REST Framework Backend & React (Vite) Frontend"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = ACCENT_ORANGE

    # Decorative Orange Line under Title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.9), Inches(1.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.fill.background()

    # 4. Content Columns (3 Columns layout)
    col_width = Inches(3.64)
    col_gap = Inches(0.4)
    start_left = Inches(0.8)
    top_pos = Inches(2.2)
    height_pos = Inches(3.8)

    columns_data = [
        {
            "title": "Architecture & Tech Stack",
            "bullets": [
                "Backend: Django 4.x + Django REST Framework",
                "Frontend: React 19 + Vite 8 + Tailwind CSS v4",
                "State & UI: shadcn/ui components, Axios client",
                "Database: SQLite for clean local storage",
                "Local Hosts: Django API (:8000), Vite Frontend (:5173)"
            ]
        },
        {
            "title": "Frontend App Core Features",
            "bullets": [
                "Home Page: Fullscreen featured video & shorts/recipes grid",
                "Recipe Explorer: Search filtering with interactive difficulty pills",
                "Recipe Details: Split hero banner, step-by-step instructions",
                "Protected Admin: Upload panel for shorts and recipes",
                "Deploy-Ready: Configured for GitHub Pages via Actions"
            ]
        },
        {
            "title": "Backend Business Logic",
            "bullets": [
                "Secure Authentication: JWT access/refresh token rotation",
                "Image Compression: Automated RGBA conversion, max 1200px JPEG q80",
                "Sliding Rate Limits: 5 short video uploads/hr, 10 recipes/day",
                "Size Validation: 50MB maximum recipe video size check",
                "Notification: Console-based auto-welcome email trigger"
            ]
        }
    ]

    for i, col in enumerate(columns_data):
        left_pos = start_left + i * (col_width + col_gap)
        
        # Add Card Background shape
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, col_width, height_pos)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        # Add Card Title
        title_box = slide.shapes.add_textbox(left_pos + Inches(0.25), top_pos + Inches(0.25), col_width - Inches(0.5), Inches(0.5))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_ct = tf_title.paragraphs[0]
        p_ct.text = col["title"]
        p_ct.font.name = "Georgia"
        p_ct.font.size = Pt(18)
        p_ct.font.bold = True
        p_ct.font.color.rgb = TEXT_DARK

        # Add Card Bullet Points
        bullets_box = slide.shapes.add_textbox(left_pos + Inches(0.25), top_pos + Inches(0.85), col_width - Inches(0.5), height_pos - Inches(1.1))
        tf_bullets = bullets_box.text_frame
        tf_bullets.word_wrap = True
        tf_bullets.margin_left = tf_bullets.margin_right = tf_bullets.margin_top = tf_bullets.margin_bottom = 0

        for idx, bullet_text in enumerate(col["bullets"]):
            p_bullet = tf_bullets.paragraphs[0] if idx == 0 else tf_bullets.add_paragraph()
            p_bullet.text = "•  " + bullet_text
            p_bullet.font.name = "Arial"
            p_bullet.font.size = Pt(11)
            p_bullet.font.color.rgb = TEXT_MUTED
            p_bullet.space_after = Pt(8)
            p_bullet.line_spacing = 1.15

    # 5. Bottom Highlights Banner
    banner_left = Inches(0.8)
    banner_top = Inches(6.25)
    banner_width = Inches(11.733)
    banner_height = Inches(0.65)

    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, banner_left, banner_top, banner_width, banner_height)
    banner.fill.solid()
    banner.fill.fore_color.rgb = TEXT_DARK
    banner.line.fill.background()

    banner_box = slide.shapes.add_textbox(banner_left, banner_top + Inches(0.18), banner_width, banner_height - Inches(0.18))
    tf_banner = banner_box.text_frame
    tf_banner.word_wrap = True
    tf_banner.margin_left = tf_banner.margin_right = tf_banner.margin_top = tf_banner.margin_bottom = 0
    p_b = tf_banner.paragraphs[0]
    p_b.text = "🔑 KEY INFRASTRUCTURE METRICS:  JWT Expiry: 30m Access / 7d Refresh  |  Max Upload size: 50MB  |  Image Resize: Auto-Scale to 1200px"
    p_b.alignment = PP_ALIGN.CENTER
    p_b.font.name = "Arial"
    p_b.font.size = Pt(10)
    p_b.font.bold = True
    p_b.font.color.rgb = BG_COLOR

    # Save presentation
    prs.save("MyFoodBlog_Summary_Slide.pptx")
    print("PowerPoint presentation 'MyFoodBlog_Summary_Slide.pptx' generated successfully!")

if __name__ == "__main__":
    create_summary_slide()
